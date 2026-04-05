from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

prs = Presentation()
slide_layout = prs.slide_layouts[1] # Title and Content
slide = prs.slides.add_slide(slide_layout)

# Apply Dark Mode Theme Space
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 20, 35)

title = slide.shapes.title
title.text = "Dark Mode Galaxy Theme"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 215, 0) # Gold
title.text_frame.paragraphs[0].font.bold = True

content = slide.placeholders[1]
content.text = "Bullet 1"
for p in content.text_frame.paragraphs:
    p.font.color.rgb = RGBColor(255, 255, 255) # White

prs.save("styled_test.pptx")
print("Saved styled_test.pptx")
