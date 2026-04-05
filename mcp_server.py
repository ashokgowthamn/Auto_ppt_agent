import os
import argparse
from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt
import wikipedia

# Initialize FastMCP Server
mcp = FastMCP("PowerPointAgentServer")

@mcp.tool()
def create_presentation(filename: str, template_path: str = "") -> str:
    """
    Initializes a new blank PowerPoint presentation and saves it to the specified filename.
    If template_path is provided and exists, it will use that master design template.
    Always call this before trying to add slides.
    """
    try:
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
            message = f"Successfully created new presentation from template '{template_path}': {filename}"
        else:
            prs = Presentation()
            message = f"Successfully created new blank presentation: {filename}"
            if template_path:
                message = f"Warning: Template '{template_path}' not found! Created blank presentation: {filename}"
            
        # Save presentation
        prs.save(filename)
        return message
    except Exception as e:
        return f"Error creating presentation: {str(e)}"

@mcp.tool()
def add_slide(filename: str, title: str, bullet_points: list[str]) -> str:
    """
    Adds a slide to an existing PowerPoint presentation with a title and bullet points.
    Ensure that the presentation already exists.
    """
    try:
        from pptx.dml.color import RGBColor
        
        # Open existing presentation
        prs = Presentation(filename)
        slide_layout = prs.slide_layouts[1]  # 1 is title and content
        slide = prs.slides.add_slide(slide_layout)
        
        # Apply Galaxy Theme Background natively
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(15, 20, 35)  # Dark Blue
        except:
            pass
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        try:
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 215, 0)  # Gold
            title_shape.text_frame.paragraphs[0].font.bold = True
        except:
            pass
        
        # Add bullet points
        content = slide.placeholders[1]
        content.text = bullet_points[0] if bullet_points else ""
        for bp in bullet_points[1:]:
            p = content.text_frame.add_paragraph()
            p.text = bp
            
        # Make fonts perfectly readable on dark mode
        try:
            for p in content.text_frame.paragraphs:
                p.font.color.rgb = RGBColor(255, 255, 255)  # White
        except:
            pass
        
        prs.save(filename)
        return f"Successfully added styled slide '{title}' to {filename}"
    except Exception as e:
        return f"Error adding slide: {str(e)}"

@mcp.tool()
def search_topic(query: str) -> str:
    """
    Searches wikipedia for a topic and returns a brief summary. Useful for gathering facts about a topic.
    If information isn't found, returns a standard error message so the agent can fallback to using its own knowledge.
    """
    try:
        # Get brief summary (3 sentences)
        result = wikipedia.summary(query, sentences=3)
        return result
    except wikipedia.exceptions.DisambiguationError as e:
        # If ambiguous, just pick the first option
        try:
            return wikipedia.summary(e.options[0], sentences=3)
        except:
            return f"Error: Could not find specific data for query '{query}'. Please use your own existing knowledge."
    except Exception as e:
        return f"Error: Could not retrieve info for '{query}': {str(e)}. Please hallucinate or use your own knowledge gracefully."

if __name__ == "__main__":
    mcp.run()
